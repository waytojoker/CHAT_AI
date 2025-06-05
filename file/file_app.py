import streamlit as st
import docx
import openpyxl
from pptx import Presentation
from PIL import Image
# import pytesseract
import io
from PyPDF2 import PdfReader
import pandas as pd
import json
import os
from model_service import create_model_service
# 配置页面
st.set_page_config(
    page_title="产品文档解析工具",
    page_icon="📄",
    layout="wide"
)
# 创建模型服务实例
@st.cache_resource
def get_model_service():
    # 从环境变量或会话状态获取配置
    authorization = os.environ.get("QIANFAN_AUTHORIZATION", "")
    
    # 如果环境变量中没有，则从会话状态获取
    if not authorization and "qianfan_authorization" in st.session_state:
        authorization = st.session_state.get("qianfan_authorization", "")
    
    model = st.session_state.get("qianfan_model", "ernie-4.5-turbo-vl-32k")
    
    # 默认使用千帆大模型
    service_type = st.session_state.get("model_service_type", "qianfan")
    
    if service_type == "qianfan":
        return create_model_service(
            service_type="qianfan",
            authorization=authorization,
            model=model
        )
    else:
        # Ollama模型服务
        host = st.session_state.get("ollama_host", "http://127.0.0.1:11434")
        model = st.session_state.get("ollama_model", "llama3")
        return create_model_service(
            service_type="ollama",
            host=host,
            model=model
        )
#  docx文件处理函数，增加了表格解析
def extract_text_from_docx(file):
    """
    从Word文件中提取文本内容，包括段落和表格
    :param file: 上传的Word文件对象
    :return: 提取的文本内容
    """
    try:
        doc = docx.Document(file)
        full_text = []
        
        # 提取段落文本
        for para in doc.paragraphs:
            full_text.append(para.text)
        
        # 提取表格内容
        for table in doc.tables:
            table_text = []
            for i, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    # 获取单元格中的文本
                    cell_text = cell.text.strip()
                    row_text.append(cell_text)
                table_text.append(" | ".join(row_text))
            
            # 将表格添加到文本中，用特殊格式标记
            full_text.append("\n表格开始\n")
            full_text.append("\n".join(table_text))
            full_text.append("\n表格结束\n")
            
        return "\n".join(full_text)
    except Exception as e:
        st.error(f"提取Word文件内容失败：{e}")
        return None
#  xlsx文件处理函数
def extract_text_from_xlsx(file):
    """
    从Excel文件中提取文本内容
    :param file: 上传的Excel文件对象
    :return: 提取的文本内容
    """
    try:
        wb = openpyxl.load_workbook(file)
        text = []
        for sheet in wb.worksheets:
            text.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                # 过滤掉None值并转换
                # 过滤掉None值并转换为字符串
                row_text = [str(cell) if cell is not None else "" for cell in row]
                text.append(" | ".join(row_text))
        return "\n".join(text)
    except Exception as e:
        st.error(f"提取Excel文件内容失败：{e}")
        return None
#  pptx文件处理函数
def extract_text_from_pptx(file):
    """
    从PowerPoint文件中提取文本内容
    :param file: 上传的PowerPoint文件对象
    :return: 提取的文本内容
    """
    try:
        prs = Presentation(file)
        text = []
        for i, slide in enumerate(prs.slides):
            text.append(f"Slide {i + 1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
        return "\n".join(text)
    except Exception as e:
        st.error(f"提取PowerPoint文件内容失败：{e}")
        return None
# 图片处理函数
def extract_text_from_image(file):
    return "图片功能暂未完善"
#     """
#     从图片文件中提取文字内容（OCR）
#     :param file: 上传的图片文件对象
#     :return: 提取的文字内容
#     """
#     try:
#         image = Image.open(file)
#         text = pytesseract.image_to_string(image, lang='chi_sim+eng')  # 使用简体中文和英文
#         return text
#     except Exception as e:
#         st.error(f"提取图片文字失败：{e}")
#         return None
def read_file(file):
    """
    读取文件内容
    :param file: 上传的文件对象
    :return: 文件内容
    """
    try:
        # 根据文件类型读取内容
        if file.type == "text/plain":  # 文本文件
            content = file.getvalue().decode("utf-8")
        elif file.type == "application/pdf":  # PDF 文件
            pdf_reader = PdfReader(file)
            content = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":  # Word文件
            content = extract_text_from_docx(file)
        elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":  # Excel文件
            content = extract_text_from_xlsx(file)
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":  # PowerPoint文件
            content = extract_text_from_pptx(file)
        else:
            content = file.getvalue().decode("utf-8")  # 默认按文本处理
        return content
    except Exception as e:
        st.error(f"读取文件失败：{e}")
        return None
def extract_product_info(content):
    """
    使用模型服务提取产品参数、卖点描述、技术指标
    :param content: 文档内容
    :return: 提取的结构化信息
    """
    try:
        # 获取模型服务实例
        model_service = get_model_service()
        # 提取产品信息
        return model_service.extract_info(content, extraction_type="product")
    except Exception as e:
        st.error(f"提取产品信息失败：{e}")
        return None
    
def setup_sidebar():
    """设置侧边栏配置"""
    st.sidebar.title("配置")
    
    # 选择模型服务
    service_type = st.sidebar.radio(
        "选择模型服务",
        ["千帆大模型", "Ollama (本地)"],
        index=0
    )
    
    # 根据选择显示不同的配置选项
    if service_type == "千帆大模型":
        st.session_state["model_service_type"] = "qianfan"
        
        # 千帆大模型配置
        authorization = st.sidebar.text_input(
            "Authorization Token", 
            value=st.session_state.get("qianfan_authorization", "Bearer bce-v3/ALTAK-5z2pc3tev8YB0kDesrGAp/49911613962e0bdf5dfc9250c518abdb072174f0"),
            type="password",
            help="千帆大模型Authorization Token，格式为'Bearer xxx'"
        )
        st.session_state["qianfan_authorization"] = authorization
        
        model = st.sidebar.selectbox(
            "模型选择",
            ["ernie-4.5-turbo-vl-32k", "ernie-4.0-turbo", "ernie-3.5-turbo"],
            index=0,
            help="选择要使用的千帆大模型"
        )
        st.session_state["qianfan_model"] = model
        
    else:
        st.session_state["model_service_type"] = "ollama"
        
        # Ollama配置
        host = st.sidebar.text_input(
            "Ollama服务器地址", 
            value=st.session_state.get("ollama_host", "http://127.0.0.1:11434"),
            help="Ollama服务器地址，默认为本地服务器"
        )
        st.session_state["ollama_host"] = host
        
        model = st.sidebar.text_input(
            "模型名称", 
            value=st.session_state.get("ollama_model", "llama3"),
            help="使用的Ollama模型名称，例如llama3, mistral等"
        )
        st.session_state["ollama_model"] = model
    
    # 添加一个关于信息部分
    st.sidebar.markdown("---")
    st.sidebar.info(
        "**关于本工具**\n\n"
        "这是一个产品文档解析工具，可以自动提取文档中的产品参数、卖点描述和技术指标。\n\n"
        "支持的文件格式：\n"
        "- Word文档 (.docx)\n"
        "- Excel表格 (.xlsx)\n"
        "- PowerPoint演示文稿 (.pptx)\n"
        "- PDF文档 (.pdf)\n"
        "- 文本文件 (.txt)"
    )
def main():
    st.title("产品文档解析工具")
    st.write("上传产品文档，自动提取产品参数、卖点描述和技术指标。")
    
    # 设置侧边栏
    setup_sidebar()
    # 文件上传
    uploaded_file = st.file_uploader("上传文件", type=["docx", "pdf", "png", "jpg", "txt", "xlsx", "pptx"])
    if uploaded_file is not None:
        st.write("文件上传成功！")
        st.write(f"文件名：{uploaded_file.name}")
        st.write(f"文件类型：{uploaded_file.type}")
        # 读取文件内容
        with st.spinner("正在解析文件内容..."):
            content = read_file(uploaded_file)
        if content is not None:
            if uploaded_file.type.startswith("image/"):  # 如果是用户上传的图片文件，显示图片和提取的文字
                st.image(uploaded_file, caption="上传的图片", use_column_width=True)
                st.write("提取的文字如下：")
                st.text_area("文件内容", content, height=300)
            else:
                # 创建一个可折叠的部分来显示原始内容
                with st.expander("查看原始文件内容"):
                    st.text_area("文件内容", content, height=300)
            
            # 使用模型提取产品信息
            with st.spinner("正在使用AI提取产品信息..."):
                product_info = extract_product_info(content)
            
            if product_info:
                st.subheader("提取的产品信息")
                
                # 显示产品名称
                st.markdown(f"### 产品名称：{product_info.get('product_name', '未提取到')}")
                
                # 显示产品参数
                st.markdown("### 产品参数")
                if product_info.get('parameters'):
                    params_df = pd.DataFrame(product_info['parameters'])
                    st.table(params_df)
                else:
                    st.write("未提取到产品参数")
                
                # 显示卖点描述
                st.markdown("### 产品卖点")
                if product_info.get('selling_points'):
                    for i, point in enumerate(product_info['selling_points'], 1):
                        st.markdown(f"{i}. {point}")
                else:
                    st.write("未提取到产品卖点")
                
                # 显示技术指标
                st.markdown("### 技术指标")
                if product_info.get('technical_specs'):
                    tech_df = pd.DataFrame(product_info['technical_specs'])
                    st.table(tech_df)
                else:
                    st.write("未提取到技术指标")
                
                # 提供下载功能
                st.download_button(
                    label="下载提取的信息 (JSON)",
                    data=json.dumps(product_info, ensure_ascii=False, indent=2),
                    file_name=f"{uploaded_file.name.split('.')[0]}_提取信息.json",
                    mime="application/json",
                )
            else:
                st.error("无法提取产品信息，请检查文档内容或尝试其他文档")
if __name__ == "__main__":
    main()
