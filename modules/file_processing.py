import streamlit as st
import ollama
import docx  # 用于处理Word文件
import openpyxl  # 用于处理Excel文件
from pptx import Presentation  # 用于处理PowerPoint文件
from PIL import Image
import pytesseract  # 用于OCR文字识别
import io

# Ollama 客户端
client = ollama.Client(host="http://127.0.0.1:11434")

#  docx文件处理函数
def extract_text_from_docx(file):
    """
    从Word文件中提取文本内容
    :param file: 上传的Word文件对象
    :return: 提取的文本内容
    """
    try:
        doc = docx.Document(file)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        st.error(f"提取Word文件内容失败：{e}")
        return None

#  xlsx文件处理函数
def extract_text_from_xlsx(file):
    """
    从Excel文件中提取文本内容
    :param file: 上传的Excel文件对象
    :return: 提取的文本内容
    """
    try:
        wb = openpyxl.load_workbook(file)
        text = ""
        for sheet in wb.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                text += "\t".join(str(cell) for cell in row) + "\n"
        return text
    except Exception as e:
        st.error(f"提取Excel文件内容失败：{e}")
        return None

#  pptx文件处理函数
def extract_text_from_pptx(file):
    """
    从PowerPoint文件中提取文本内容
    :param file: 上传的PowerPoint文件对象
    :return: 提取的文本内容
    """
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            text += f"Slide {prs.slides.index(slide) + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"提取PowerPoint文件内容失败：{e}")
        return None

# 图片处理函数
def extract_text_from_image(file):
    return "图片功能暂未完善"
#     """
#     从图片文件中提取文字内容（OCR）
#     :param file: 上传的图片文件对象
#     :return: 提取的文字内容
#     """
#     try:
#         image = Image.open(file)
#         text = pytesseract.image_to_string(image, lang='chi_sim+eng')  # 使用简体中文和英文
#         return text
#     except Exception as e:
#         st.error(f"提取图片文字失败：{e}")
#         return None

def read_file(file):
    """
    读取文件内容
    :param file: 上传的文件对象
    :return: 文件内容
    """
    try:
        # 根据文件类型读取内容
        if file.type == "text/plain":  # 文本文件
            content = file.getvalue().decode("utf-8")
        elif file.type == "application/pdf":  # PDF 文件
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(file)
            content = "\n".join([page.extract_text() for page in pdf_reader.pages])
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":  # Word文件
            content = extract_text_from_docx(file)
        elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":  # Excel文件
            content = extract_text_from_xlsx(file)
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":  # PowerPoint文件
            content = extract_text_from_pptx(file)

        else:
            content = file.getvalue().decode("utf-8")  # 默认按文本处理
        return content
    except Exception as e:
        st.error(f"读取文件失败：{e}")
        return None


def get_file_content(uploaded_files):
    if uploaded_files is None:
        return None

    # 若传入的是单个文件对象，将其转换为列表
    if not isinstance(uploaded_files, list):
        uploaded_files = [uploaded_files]

    file_content = ""
    for uploaded_file in uploaded_files:
        # 读取文件内容
        file_content += f"文件名：{uploaded_file.name}\n"
        file_content += f"文件类型：{uploaded_file.type}\n"
        content = read_file(uploaded_file)
        if content:
            file_content += content + "\n\n"  # 添加空行分隔不同文件

    return file_content

def main():
    st.title("文件上传测试")
    st.write("上传一个文件，我将读取并显示其内容。")

    # 文件上传
    uploaded_file = st.file_uploader("上传文件", type=["docx", "pdf", "png", "jpg", "txt", "xlsx", "pptx"])

    if uploaded_file is not None:
        st.write("文件上传成功！")
        st.write(f"文件名：{uploaded_file.name}")
        st.write(f"文件类型：{uploaded_file.type}")

        # 读取文件内容
        content = read_file(uploaded_file)
        st.write("文件内容如下：")
        st.text(content)


if __name__ == "__main__":
    main()